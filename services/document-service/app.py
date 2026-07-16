from __future__ import annotations

import io
import os
import subprocess
import tempfile
import zipfile
from pathlib import Path
from typing import Literal
from urllib.parse import quote
from xml.etree import ElementTree
from xml.sax.saxutils import escape

from docling.document_converter import DocumentConverter
from docx import Document
from docx.enum.table import WD_TABLE_ALIGNMENT
from fastapi import FastAPI, File, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from pydantic import BaseModel, Field


SUPPORTED_INPUTS = {".pdf", ".docx", ".pptx", ".xlsx", ".html", ".md", ".txt", ".csv", ".hwpx"}
TEMPLATE_DIR = Path(os.getenv("HERION_TEMPLATE_DIR", "/app/templates"))

app = FastAPI(title="Herian Document Service", version="0.1.0")
allowed_origins = [origin.strip() for origin in os.getenv(
    "HERION_ALLOWED_ORIGINS", "https://icanjt7.github.io,http://localhost:3000"
).split(",") if origin.strip()]
app.add_middleware(
    CORSMiddleware,
    allow_origins=allowed_origins,
    allow_credentials=False,
    allow_methods=["GET", "POST", "OPTIONS"],
    allow_headers=["*"],
    expose_headers=["Content-Disposition"],
)

converter = DocumentConverter()


class TableSpec(BaseModel):
    title: str = ""
    headers: list[str] = Field(default_factory=list)
    rows: list[list[str]] = Field(default_factory=list)


class SectionSpec(BaseModel):
    heading: str
    paragraphs: list[str] = Field(default_factory=list)
    bullets: list[str] = Field(default_factory=list)
    tables: list[TableSpec] = Field(default_factory=list)


class ReportRequest(BaseModel):
    title: str
    subtitle: str = ""
    author: str = ""
    format: Literal["docx", "pdf", "pptx", "hwpx"] = "docx"
    template_id: str = ""
    sections: list[SectionSpec] = Field(default_factory=list)


def safe_filename(value: str, fallback: str = "herian-report") -> str:
    cleaned = "".join(char for char in value if char.isalnum() or char in " -_().").strip(" .")
    return cleaned[:80] or fallback


def template_path(template_id: str, suffix: str) -> Path | None:
    if not template_id:
        return None
    safe_id = safe_filename(template_id, "")
    candidate = (TEMPLATE_DIR / f"{safe_id}{suffix}").resolve()
    if TEMPLATE_DIR.resolve() not in candidate.parents or not candidate.is_file():
        raise HTTPException(422, f"등록된 {suffix} 템플릿을 찾을 수 없습니다: {template_id}")
    return candidate


def parse_hwpx(data: bytes) -> tuple[str, dict]:
    with zipfile.ZipFile(io.BytesIO(data)) as package:
        section_names = sorted(
            name for name in package.namelist()
            if name.lower().startswith("contents/section") and name.lower().endswith(".xml")
        )
        if not section_names:
            raise HTTPException(422, "HWPX 본문 section XML을 찾지 못했습니다.")
        sections = []
        markdown_parts = []
        for index, name in enumerate(section_names, 1):
            root = ElementTree.fromstring(package.read(name))
            paragraphs = []
            for paragraph in (node for node in root.iter() if node.tag.rsplit("}", 1)[-1] == "p"):
                text = "".join(
                    node.text or "" for node in paragraph.iter()
                    if node.tag.rsplit("}", 1)[-1] == "t"
                ).strip()
                if text:
                    paragraphs.append(text)
            sections.append({"index": index, "source": name, "paragraphs": paragraphs})
            markdown_parts.append(f"## HWPX 구역 {index}\n\n" + "\n\n".join(paragraphs))
    return "\n\n".join(markdown_parts), {"format": "hwpx", "sections": sections}


@app.get("/health")
def health() -> dict:
    return {"status": "ok", "service": "herian-document-service"}


@app.post("/v1/documents/parse")
async def parse_document(file: UploadFile = File(...)) -> dict:
    suffix = Path(file.filename or "").suffix.lower()
    if suffix not in SUPPORTED_INPUTS:
        raise HTTPException(415, f"지원하지 않는 형식입니다: {suffix}")
    data = await file.read()

    if suffix == ".hwpx":
        markdown, document = parse_hwpx(data)
        return {"filename": file.filename, "parser": "herian-hwpx", "markdown": markdown, "document": document}

    with tempfile.TemporaryDirectory(prefix="herion-parse-") as temp_dir:
        source = Path(temp_dir) / safe_filename(file.filename or f"document{suffix}")
        source.write_bytes(data)
        try:
            result = converter.convert(source)
            document = result.document
            return {
                "filename": file.filename,
                "parser": "docling",
                "markdown": document.export_to_markdown(),
                "document": document.export_to_dict(),
            }
        except Exception as error:
            raise HTTPException(422, f"문서 구조를 읽지 못했습니다: {error}") from error


def build_docx(request: ReportRequest, target: Path) -> None:
    template = template_path(request.template_id, ".docx")
    document = Document(str(template)) if template else Document()
    document.add_heading(request.title, 0)
    if request.subtitle:
        document.add_paragraph(request.subtitle)
    if request.author:
        document.add_paragraph(request.author)
    for section in request.sections:
        document.add_heading(section.heading, level=1)
        for paragraph in section.paragraphs:
            document.add_paragraph(paragraph)
        for bullet in section.bullets:
            document.add_paragraph(bullet, style="List Bullet")
        for table_spec in section.tables:
            if table_spec.title:
                title_paragraph = document.add_paragraph()
                title_paragraph.add_run(table_spec.title).bold = True
            columns = max(len(table_spec.headers), max((len(row) for row in table_spec.rows), default=0))
            if not columns:
                continue
            table = document.add_table(rows=1 if table_spec.headers else 0, cols=columns)
            table.style = "Table Grid"
            table.alignment = WD_TABLE_ALIGNMENT.CENTER
            if table_spec.headers:
                for index, value in enumerate(table_spec.headers):
                    table.rows[0].cells[index].text = value
            for values in table_spec.rows:
                cells = table.add_row().cells
                for index, value in enumerate(values[:columns]):
                    cells[index].text = value
    document.save(target)


def build_pptx(request: ReportRequest, target: Path) -> None:
    template = template_path(request.template_id, ".pptx")
    presentation = Presentation(str(template)) if template else Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = request.title
    title_slide.placeholders[1].text = "\n".join(part for part in [request.subtitle, request.author] if part)
    for section in request.sections:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = section.heading
        frame = slide.placeholders[1].text_frame
        frame.clear()
        items = [*section.paragraphs, *section.bullets]
        for index, value in enumerate(items[:12]):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = value
            paragraph.font.size = Pt(18)
        for table_spec in section.tables[:1]:
            rows = 1 + len(table_spec.rows)
            columns = max(len(table_spec.headers), max((len(row) for row in table_spec.rows), default=0))
            if not columns:
                continue
            shape = slide.shapes.add_table(rows, columns, Inches(0.6), Inches(3.6), Inches(12.1), Inches(3.1))
            table = shape.table
            for column, value in enumerate(table_spec.headers[:columns]):
                table.cell(0, column).text = value
            for row_index, values in enumerate(table_spec.rows, 1):
                for column, value in enumerate(values[:columns]):
                    table.cell(row_index, column).text = value
    presentation.save(target)


def build_hwpx_from_template(request: ReportRequest, target: Path) -> None:
    template = template_path(request.template_id, ".hwpx")
    if not template:
        raise HTTPException(422, "HWPX 생성에는 승인된 기관 template_id가 필요합니다.")
    body = "\n\n".join(
        [request.subtitle, *(
            f"{section.heading}\n" + "\n".join([*section.paragraphs, *section.bullets])
            for section in request.sections
        )]
    ).strip()
    replacements = {
        "{{TITLE}}": escape(request.title),
        "{{SUBTITLE}}": escape(request.subtitle),
        "{{AUTHOR}}": escape(request.author),
        "{{BODY}}": escape(body),
    }
    with zipfile.ZipFile(template) as source, zipfile.ZipFile(target, "w") as output:
        for item in source.infolist():
            payload = source.read(item.filename)
            if item.filename.lower().startswith("contents/section") and item.filename.lower().endswith(".xml"):
                text = payload.decode("utf-8")
                for marker, value in replacements.items():
                    text = text.replace(marker, value)
                payload = text.encode("utf-8")
            output.writestr(item, payload)


@app.post("/v1/reports")
def create_report(request: ReportRequest) -> StreamingResponse:
    with tempfile.TemporaryDirectory(prefix="herion-report-") as temp_dir:
        directory = Path(temp_dir)
        stem = safe_filename(request.title)
        if request.format in {"docx", "pdf"}:
            docx_path = directory / f"{stem}.docx"
            build_docx(request, docx_path)
            if request.format == "pdf":
                completed = subprocess.run(
                    ["soffice", "--headless", "--convert-to", "pdf", "--outdir", str(directory), str(docx_path)],
                    capture_output=True, text=True, timeout=90, check=False,
                )
                target = directory / f"{stem}.pdf"
                if completed.returncode or not target.exists():
                    raise HTTPException(500, f"PDF 변환에 실패했습니다: {completed.stderr.strip()}")
                media_type = "application/pdf"
            else:
                target = docx_path
                media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        elif request.format == "pptx":
            target = directory / f"{stem}.pptx"
            build_pptx(request, target)
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        else:
            target = directory / f"{stem}.hwpx"
            build_hwpx_from_template(request, target)
            media_type = "application/hwp+zip"

        payload = target.read_bytes()
        headers = {
            "Content-Disposition": (
                f'attachment; filename="herian-report.{request.format}"; '
                f"filename*=UTF-8''{quote(target.name)}"
            )
        }
        return StreamingResponse(io.BytesIO(payload), media_type=media_type, headers=headers)
